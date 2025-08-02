# # Updated backend/services/preprocessing.py

# import re
# from pathlib import Path
# import mimetypes

# def extract_text_from_pdf(path: Path) -> str:
#     import fitz  # PyMuPDF
#     doc = fitz.open(str(path))
#     text = ""
#     for page in doc:
#         text += page.get_text()
#     return text

# def extract_text_from_docx(path: Path) -> str:
#     from docx import Document
#     doc = Document(path)
#     full_text = []
#     for para in doc.paragraphs:
#         full_text.append(para.text)
#     return "\n".join(full_text)

# def extract_text_from_pptx(path: Path) -> str:
#     from pptx import Presentation
#     prs = Presentation(path)
#     text = []
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text.append(shape.text)
#     return "\n".join(text)

# def extract_text(path: Path) -> str:
#     suffix = path.suffix.lower()
#     if suffix == ".txt":
#         return path.read_text(encoding="utf-8")
#     elif suffix == ".pdf":
#         return extract_text_from_pdf(path)
#     elif suffix == ".docx":
#         return extract_text_from_docx(path)
#     elif suffix == ".pptx":
#         return extract_text_from_pptx(path)
#     else:
#         raise ValueError(f"Unsupported file type: {suffix}")

# def clean_text(text: str) -> str:
#     cleaned = re.sub(r"\s+", " ", text).strip()
#     return cleaned

# def clean_and_save_document(file_path: Path) -> Path:
#     raw_text = extract_text(file_path)
#     cleaned = clean_text(raw_text)
#     cleaned_file = file_path.with_suffix(".cleaned.txt")
#     with open(cleaned_file, "w", encoding="utf-8") as f:
#         f.write(cleaned)
#     return cleaned_file


import re
from pathlib import Path

def extract_text_from_pdf(path: Path) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(str(path))
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(para.text for para in doc.paragraphs)

def extract_text_from_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".txt":
        return path.read_text(encoding="utf-8")
    elif suffix == ".pdf":
        return extract_text_from_pdf(path)
    elif suffix == ".docx":
        return extract_text_from_docx(path)
    elif suffix == ".pptx":
        return extract_text_from_pptx(path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")

def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()

def clean_and_save_document(file_path: Path) -> Path:
    raw_text = extract_text(file_path)
    cleaned = clean_text(raw_text)
    cleaned_file = file_path.with_suffix(".cleaned.txt")
    cleaned_file.write_text(cleaned, encoding="utf-8")
    return cleaned_file
